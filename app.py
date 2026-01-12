import requests
from pptx import Presentation
import tiktoken


INPUT_PPTX = "input.pptx"               # Path to input PPT
TARGET_LANGUAGE = "French"               # Target language
OUTPUT_PPTX = "output_french.pptx"       # Output PPT

TOKEN_LIMIT = 3500
MAX_RETRIES = 3

AZURE_ENDPOINT = "AZURE_OPENAI_ENDPOINT"
DEPLOYMENT_NAME = "AZURE_OPENAI_MODEL_NAME"
API_VERSION = "AZURE_OPENAI_PREVIEW_API_VERSION"
AZURE_KEY = "AZURE_OPENAI_KEY"

HEADERS = {
    "Content-Type": "application/json",
    "api-key": AZURE_KEY
}


# -------------------- TOKEN UTILS --------------------

def estimate_token_count(text):
    encoding = tiktoken.get_encoding("cl100k_base")
    return len(encoding.encode(text))


# -------------------- TRANSLATION --------------------

def translate_text_batch(texts, target_language):
    texts = [t.strip() for t in texts if t.strip()]
    if not texts:
        return texts

    combined_text = "\n\n".join(texts)

    if estimate_token_count(combined_text) > TOKEN_LIMIT:
        return None

    api_url = (
        f"{AZURE_ENDPOINT}/openai/deployments/"
        f"{DEPLOYMENT_NAME}/chat/completions?api-version={API_VERSION}"
    )

    body = {
        "messages": [
            {
                "role": "system",
                "content": (
                    f"You are a professional translator specializing in {target_language}. "
                    "Translate the following text very accurately. "
                    "Do not alter numbers, formatting, or structure."
                ),
            },
            {
                "role": "user",
                "content": f"Translate this into {target_language}:\n{combined_text}",
            }
        ],
        "temperature": 0,
        "max_tokens": 4000,
    }

    for attempt in range(MAX_RETRIES):
        try:
            response = requests.post(api_url, headers=HEADERS, json=body, timeout=30)
            data = response.json()
            if response.status_code == 200:
                return data["choices"][0]["message"]["content"].split("\n\n")
            else:
                print("API Error:", data)
        except Exception as e:
            print(f"Retry {attempt + 1} failed:", e)

    return None

# -------------------- FORMATTING SAFE REPLACE --------------------

def replace_text_in_ref(obj, new_text):
    def get_font_color(font):
        try:
            if font.color and font.color.rgb:
                return font.color.rgb
        except:
            return None

    if hasattr(obj, "text_frame"):
        tf = obj.text_frame
        while len(tf.paragraphs) > 1:
            tf._element.remove(tf.paragraphs[-1]._element)

        para = tf.paragraphs[0]
    else:
        para = obj

    if para.runs:
        run = para.runs[0]
        font = run.font
        size, name = font.size, font.name
        bold, italic = font.bold, font.italic
        color = get_font_color(font)
    else:
        size = name = bold = italic = color = None

    level = para.level
    para.clear()

    new_run = para.add_run()
    new_run.text = new_text
    font = new_run.font

    if size:
        font.size = size
    if name:
        font.name = name
    if bold is not None:
        font.bold = bold
    if italic is not None:
        font.italic = italic
    if color:
        font.color.rgb = color

    para.level = level

# -------------------- TEXT COLLECTION --------------------

def collect_text_items_from_shape(shape, items):
    if hasattr(shape, "shapes"):
        for sub in shape.shapes:
            collect_text_items_from_shape(sub, items)

    elif hasattr(shape, "text_frame") and shape.text_frame:
        for para in shape.text_frame.paragraphs:
            text = " ".join(run.text for run in para.runs).strip()
            if text:
                items.append((text, para))

    elif hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    items.append((cell.text.strip(), cell))

# -------------------- PPT TRANSLATION --------------------

def translate_pptx(input_pptx, output_pptx, target_language):
    prs = Presentation(input_pptx)

    for slide in prs.slides:
        text_items = []

        for shape in slide.shapes:
            collect_text_items_from_shape(shape, text_items)

        batch, token_count = [], 0

        for text, ref in text_items:
            tokens = estimate_token_count(text)
            if token_count + tokens > TOKEN_LIMIT:
                translated = translate_text_batch([t for t, _ in batch], target_language)
                if translated:
                    for (_, obj), t in zip(batch, translated):
                        replace_text_in_ref(obj, t)
                batch, token_count = [], 0

            batch.append((text, ref))
            token_count += tokens

        if batch:
            translated = translate_text_batch([t for t, _ in batch], target_language)
            if translated:
                for (_, obj), t in zip(batch, translated):
                    replace_text_in_ref(obj, t)

    prs.save(output_pptx)

# -------------------- MAIN --------------------

if __name__ == "__main__":
    print("🔄 Translating PowerPoint...")
    translate_pptx(INPUT_PPTX, OUTPUT_PPTX, TARGET_LANGUAGE)
    print(f"✅ Translation complete. Saved as: {OUTPUT_PPTX}")
